from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path

import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches


def _fig_to_png_bytes(fig) -> bytes:
    # Requires kaleido
    return pio.to_image(fig, format="png", scale=2)


def build_pptx(output_path: str | Path, figures: Iterable, title: str = "Client Report") -> Path:
    output_path = Path(output_path)
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    # One slide per figure
    for fig in figures:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = (
            getattr(fig.layout, "title", {}).text if hasattr(fig, "layout") else "Chart"
        )

        png = _fig_to_png_bytes(fig)
        img_path = output_path.with_suffix(".tmp.png")
        img_path.write_bytes(png)
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1.5), width=Inches(8))
        try:
            img_path.unlink()
        except OSError:
            pass

    prs.save(str(output_path))
    return output_path
